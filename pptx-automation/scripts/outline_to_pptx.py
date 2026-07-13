#!/usr/bin/env python3
"""Create an editable PPTX and HTML preview from a structured deck outline."""

from __future__ import annotations

import argparse
import html
import json
import shutil
import subprocess
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.58
DEFAULT_BRAND_PROFILE = Path(__file__).resolve().parents[1] / "assets" / "brand-profile.json"


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def load_brand(path: Path | None) -> dict[str, Any]:
    brand_path = path or DEFAULT_BRAND_PROFILE
    if brand_path.exists():
        return load_json(brand_path)
    return {}


def parse_brand_override(raw: str) -> tuple[str, str]:
    if "=" not in raw:
        raise argparse.ArgumentTypeError("Brand override must use key=value format.")
    key, value = raw.split("=", 1)
    key = key.strip()
    if not key:
        raise argparse.ArgumentTypeError("Brand override key cannot be empty.")
    return key, value.strip()


def merge_brand(data: dict[str, Any], profile: dict[str, Any], overrides: list[tuple[str, str]] | None) -> dict[str, Any]:
    merged = dict(profile)
    outline_brand = data.get("brand") or {}
    merged.update(outline_brand)
    override_keys = set()
    for key, value in overrides or []:
        merged[key] = value
        override_keys.add(key)
    if "name" in outline_brand and "footer_text" not in outline_brand and "footer_text" not in override_keys:
        merged["footer_text"] = outline_brand["name"]
    if "name" in override_keys and "footer_text" not in override_keys:
        merged["footer_text"] = merged["name"]
    data["brand"] = merged
    return merged


def parse_color(value: str | None, fallback: str) -> RGBColor:
    raw = (value or fallback).strip().lstrip("#")
    if len(raw) != 6:
        raw = fallback.lstrip("#")
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def add_text(shape, text: str, size: int, bold: bool = False, color: RGBColor | None = None):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


def add_title(slide, title: str, subtitle: str | None, primary: RGBColor):
    title_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(0.42), Inches(9.9), Inches(0.75))
    add_text(title_box, title, 28, True, primary)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(1.12), Inches(10.8), Inches(0.38))
        add_text(sub_box, subtitle, 12, False, RGBColor(85, 95, 110))


def add_footer(slide, index: int, primary: RGBColor, footer_text: str = ""):
    line = slide.shapes.add_shape(1, Inches(MARGIN), Inches(7.02), Inches(12.15), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 225, 232)
    line.line.fill.background()
    if footer_text:
        footer = slide.shapes.add_textbox(Inches(MARGIN), Inches(6.86), Inches(6.4), Inches(0.25))
        add_text(footer, footer_text, 8, False, RGBColor(120, 132, 150))
    page = slide.shapes.add_textbox(Inches(12.18), Inches(6.86), Inches(0.55), Inches(0.25))
    tf = add_text(page, f"{index:02d}", 8, True, primary)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_cover(prs: Presentation, slide_data: dict[str, Any], brand: dict[str, Any], primary: RGBColor, secondary: RGBColor):
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = primary
    bg.line.fill.background()
    accent = slide.shapes.add_shape(1, Inches(0), Inches(6.9), Inches(SLIDE_W), Inches(0.6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = secondary
    accent.line.fill.background()
    title = slide.shapes.add_textbox(Inches(0.85), Inches(2.4), Inches(9.5), Inches(1.0))
    add_text(title, slide_data.get("title", ""), 38, True, RGBColor(255, 255, 255))
    subtitle = slide_data.get("subtitle")
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.88), Inches(3.45), Inches(8.8), Inches(0.55))
        add_text(sub, subtitle, 16, False, RGBColor(235, 240, 246))
    meta = slide_data.get("meta") or slide_data.get("company") or brand.get("name") or ""
    if meta:
        box = slide.shapes.add_textbox(Inches(0.9), Inches(5.72), Inches(6.5), Inches(0.35))
        add_text(box, str(meta), 11, False, RGBColor(235, 240, 246))


def add_section(prs: Presentation, slide_data: dict[str, Any], primary: RGBColor, secondary: RGBColor, index: int):
    slide = blank_slide(prs)
    left = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(3.0), Inches(SLIDE_H))
    left.fill.solid()
    left.fill.fore_color.rgb = primary
    left.line.fill.background()
    tag = slide.shapes.add_textbox(Inches(0.72), Inches(2.7), Inches(1.6), Inches(0.35))
    add_text(tag, f"{index:02d}", 18, True, RGBColor(255, 255, 255))
    title = slide.shapes.add_textbox(Inches(3.7), Inches(2.72), Inches(8.2), Inches(0.9))
    add_text(title, slide_data.get("title", ""), 32, True, primary)
    if slide_data.get("subtitle"):
        sub = slide.shapes.add_textbox(Inches(3.72), Inches(3.7), Inches(7.4), Inches(0.42))
        add_text(sub, slide_data["subtitle"], 14, False, RGBColor(85, 95, 110))
    rule = slide.shapes.add_shape(1, Inches(3.72), Inches(4.35), Inches(2.2), Inches(0.07))
    rule.fill.solid()
    rule.fill.fore_color.rgb = secondary
    rule.line.fill.background()


def add_bullets(slide, bullets: list[Any], x: float, y: float, w: float, h: float):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.level = 0
        p.font.size = Pt(16 if len(bullets) <= 4 else 13)
        p.font.color.rgb = RGBColor(42, 49, 60)
        p.space_after = Pt(8)


def add_content(prs: Presentation, slide_data: dict[str, Any], brand: dict[str, Any], primary: RGBColor, index: int):
    slide = blank_slide(prs)
    add_title(slide, slide_data.get("title", ""), slide_data.get("subtitle"), primary)
    body = slide_data.get("body")
    bullets = slide_data.get("bullets") or ([] if body is None else [body])
    add_bullets(slide, bullets, MARGIN + 0.1, 1.85, 11.6, 4.7)
    add_footer(slide, index, primary, str(brand.get("footer_text") or brand.get("name") or ""))


def add_cards(prs: Presentation, slide_data: dict[str, Any], brand: dict[str, Any], primary: RGBColor, secondary: RGBColor, index: int):
    slide = blank_slide(prs)
    add_title(slide, slide_data.get("title", ""), slide_data.get("subtitle"), primary)
    cards = slide_data.get("cards") or []
    cols = 3 if len(cards) > 2 else max(1, len(cards))
    card_w = (11.9 - (cols - 1) * 0.25) / cols
    for i, card in enumerate(cards[:6]):
        row = i // cols
        col = i % cols
        x = MARGIN + col * (card_w + 0.25)
        y = 1.85 + row * 2.25
        rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(card_w), Inches(1.9))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(248, 250, 252)
        rect.line.color.rgb = RGBColor(218, 225, 235)
        title = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.22), Inches(card_w - 0.4), Inches(0.34))
        add_text(title, str(card.get("title", "")), 15, True, primary)
        body = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.72), Inches(card_w - 0.4), Inches(0.95))
        add_text(body, str(card.get("body", "")), 11, False, RGBColor(55, 65, 81))
    add_footer(slide, index, primary, str(brand.get("footer_text") or brand.get("name") or ""))


def add_table(prs: Presentation, slide_data: dict[str, Any], brand: dict[str, Any], primary: RGBColor, index: int):
    slide = blank_slide(prs)
    add_title(slide, slide_data.get("title", ""), slide_data.get("subtitle"), primary)
    table_data = slide_data.get("table") or {}
    headers = table_data.get("headers") or []
    rows = table_data.get("rows") or []
    cols = max(1, len(headers))
    table_shape = slide.shapes.add_table(len(rows) + 1, cols, Inches(MARGIN), Inches(1.8), Inches(12.15), Inches(4.7))
    table = table_shape.table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.font.size = Pt(10)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row[:cols]):
            cell = table.cell(r, c)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(42, 49, 60)
    add_footer(slide, index, primary, str(brand.get("footer_text") or brand.get("name") or ""))


def add_closing(prs: Presentation, slide_data: dict[str, Any], brand: dict[str, Any], primary: RGBColor, secondary: RGBColor, index: int):
    slide = blank_slide(prs)
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(248, 250, 252)
    band.line.fill.background()
    title = slide.shapes.add_textbox(Inches(1.1), Inches(2.65), Inches(9.6), Inches(0.8))
    add_text(title, slide_data.get("title", "Thank you"), 34, True, primary)
    if slide_data.get("subtitle"):
        sub = slide.shapes.add_textbox(Inches(1.12), Inches(3.55), Inches(8.4), Inches(0.45))
        add_text(sub, slide_data["subtitle"], 15, False, RGBColor(55, 65, 81))
    contact = slide_data.get("contact") or " | ".join(
        item for item in [brand.get("email", ""), brand.get("website", ""), brand.get("phone", "")] if item
    )
    if contact:
        contact_box = slide.shapes.add_textbox(Inches(1.12), Inches(4.08), Inches(8.4), Inches(0.35))
        add_text(contact_box, str(contact), 10, False, RGBColor(85, 95, 110))
    rule = slide.shapes.add_shape(1, Inches(1.12), Inches(4.38), Inches(2.2), Inches(0.08))
    rule.fill.solid()
    rule.fill.fore_color.rgb = secondary
    rule.line.fill.background()
    add_footer(slide, index, primary, str(brand.get("footer_text") or brand.get("name") or ""))


def build_pptx(data: dict[str, Any], out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    brand = data.get("brand") or {}
    primary = parse_color(brand.get("primary"), "#1F4E79")
    secondary = parse_color(brand.get("secondary"), "#22A699")
    slides = data.get("slides") or []
    if not slides:
        slides = [{"type": "cover", "title": data.get("title", "Presentation"), "subtitle": data.get("subtitle", "")}]
    for index, slide_data in enumerate(slides, start=1):
        kind = (slide_data.get("type") or "content").lower()
        if kind == "cover":
            add_cover(prs, slide_data, brand, primary, secondary)
        elif kind == "section":
            add_section(prs, slide_data, primary, secondary, index)
        elif kind == "cards":
            add_cards(prs, slide_data, brand, primary, secondary, index)
        elif kind == "table":
            add_table(prs, slide_data, brand, primary, index)
        elif kind == "closing":
            add_closing(prs, slide_data, brand, primary, secondary, index)
        else:
            add_content(prs, slide_data, brand, primary, index)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def write_html(data: dict[str, Any], out_path: Path):
    brand = data.get("brand") or {}
    primary = html.escape(brand.get("primary") or "#1F4E79")
    secondary = html.escape(brand.get("secondary") or "#22A699")
    font_family = html.escape(brand.get("font_family") or "-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif")
    parts = [
        "<!doctype html><html><head><meta charset='utf-8'>",
        "<meta name='viewport' content='width=device-width, initial-scale=1'>",
        "<style>",
        f"body{{margin:0;background:#eef2f7;font-family:{font_family};color:#263241}}",
        ".deck{max-width:1120px;margin:28px auto;padding:0 20px}.slide{aspect-ratio:16/9;background:white;margin:0 0 24px;padding:48px 56px;box-shadow:0 12px 34px rgba(15,23,42,.12);box-sizing:border-box}",
        f".cover{{background:{primary};color:white}}.section{{border-left:44px solid {primary}}}",
        f"h1{{margin:0 0 14px;font-size:34px;color:{primary}}}.cover h1{{color:white}}h2{{font-size:18px;color:#647083;margin:0 0 26px}}",
        "ul{font-size:22px;line-height:1.55}.cards{display:grid;grid-template-columns:repeat(3,1fr);gap:18px}.card{border:1px solid #dce4ef;background:#f8fafc;padding:20px}",
        f".card h3{{margin:0 0 10px;color:{primary}}}.accent{{width:120px;height:7px;background:{secondary};margin-top:24px}}",
        "table{width:100%;border-collapse:collapse;font-size:16px}th{background:#1f4e79;color:white}td,th{border:1px solid #d8e0eb;padding:10px;text-align:left;vertical-align:top}",
        "</style></head><body><main class='deck'>",
    ]
    for slide in data.get("slides", []):
        kind = html.escape(slide.get("type") or "content")
        parts.append(f"<section class='slide {kind}'>")
        parts.append(f"<h1>{html.escape(str(slide.get('title','')))}</h1>")
        if slide.get("subtitle"):
            parts.append(f"<h2>{html.escape(str(slide['subtitle']))}</h2>")
        if slide.get("bullets"):
            parts.append("<ul>")
            parts.extend(f"<li>{html.escape(str(item))}</li>" for item in slide["bullets"])
            parts.append("</ul>")
        if slide.get("body"):
            parts.append(f"<p>{html.escape(str(slide['body']))}</p>")
        if slide.get("cards"):
            parts.append("<div class='cards'>")
            for card in slide["cards"]:
                parts.append("<article class='card'>")
                parts.append(f"<h3>{html.escape(str(card.get('title','')))}</h3>")
                parts.append(f"<p>{html.escape(str(card.get('body','')))}</p>")
                parts.append("</article>")
            parts.append("</div>")
        if slide.get("table"):
            table = slide["table"]
            parts.append("<table><thead><tr>")
            parts.extend(f"<th>{html.escape(str(h))}</th>" for h in table.get("headers", []))
            parts.append("</tr></thead><tbody>")
            for row in table.get("rows", []):
                parts.append("<tr>")
                parts.extend(f"<td>{html.escape(str(c))}</td>" for c in row)
                parts.append("</tr>")
            parts.append("</tbody></table>")
        parts.append("<div class='accent'></div></section>")
    parts.append("</main></body></html>")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("\n".join(parts), encoding="utf-8")


def convert_pdf(pptx_path: Path):
    converter = shutil.which("soffice") or shutil.which("libreoffice")
    if not converter:
        raise RuntimeError("LibreOffice/soffice is not available for PDF conversion.")
    subprocess.run(
        [converter, "--headless", "--convert-to", "pdf", "--outdir", str(pptx_path.parent), str(pptx_path)],
        check=True,
    )


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("outline", type=Path, help="Path to deck outline JSON")
    parser.add_argument("--out", type=Path, required=True, help="Output PPTX path")
    parser.add_argument("--html", type=Path, help="Optional HTML preview path")
    parser.add_argument("--pdf", action="store_true", help="Also export PDF using LibreOffice/soffice")
    parser.add_argument("--brand-profile", type=Path, help="Optional brand profile JSON. Defaults to assets/brand-profile.json when present.")
    parser.add_argument("--brand", action="append", type=parse_brand_override, help="One-time brand override in key=value format. May be repeated.")
    args = parser.parse_args()

    data = load_json(args.outline)
    merge_brand(data, load_brand(args.brand_profile), args.brand)
    build_pptx(data, args.out)
    if args.html:
        write_html(data, args.html)
    if args.pdf:
        convert_pdf(args.out)
    print(f"Wrote {args.out}")
    if args.html:
        print(f"Wrote {args.html}")


if __name__ == "__main__":
    main()
